from pptx import Presentation

# Abre o seu arquivo local
prs = Presentation("NamingRights_IBMEC_Layout_Profissional.pptx")

# Dicionário de tradução para caracteres matemáticos limpos
unicode_formulas = {
    "Y_real,t = (Y_nominal,t / IGP-M_t) × 100": "Yᵣₑₐ₁,ₜ = ( Yₙₒₘᵢₙₐ₁,ₜ / IGP-Mₜ ) × 100",
    "F_Chow = [(RSS_R − (RSS_pre + RSS_pos)) / k] / [(RSS_pre + RSS_pos) / (N − 2k)]": "F_Chow = [ (RSS_R − (RSS_ₚᵣₑ + RSS_ₚₒₛ)) / k ] / [ (RSS_¼ᵣₑ + RSS_ₚₒₛ) / (N − 2k) ]",
    "Y_t = β₀ + β₁·Trend_t + β₂·Dummy_t + β₃·(Dummy_t × Trend_t) + ε_t": "Yₜ = β₀ + β₁·Trendₜ + β₂·Dummyₜ + β₃·(Dummyₜ × Trendₜ) + εₜ",
    "Y_it = α₀ + α₁·Tratado_i + α₂·Tempo_t + α₃·(Tratado_i × Tempo_t) + μ_it": "Yᵢₜ = α₀ + α₁·Tratadoᵢ + α₂·Tempoₜ + α₃·(Tratadoᵢ × Tempoₜ) + μᵢₜ",
    "α₃ = (Ȳ_Tratado,pós − Ȳ_Tratado,pré) − (Ȳ_Controle,pós − Ȳ_Controle,pré)": "α₃ = (Ȳ_Tratado,ₚₒₛ − Ȳ_Tratado,ₚᵣₑ) − (Ȳ_Controle,ₚₒₛ − Ȳ_Controle,ₚᵣₑ)"
}

# Varre os slides substituindo o texto desconfigurado pelas fórmulas limpas
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for orig, new in unicode_formulas.items():
                    if orig in p.text:
                        p.text = p.text.replace(orig, new)
                        p.font.name = 'Arial'

# Salva por cima com as correções prontas
prs.save("NamingRights_IBMEC_Layout_Profissional.pptx")
print("Pronto! Suas fórmulas foram corrigidas direto no seu arquivo.")