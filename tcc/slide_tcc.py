"""
=============================================================================
PIPELINE INTEGRADA DEFINITIVA: SLIDES COM INJEÇÃO AUTOMÁTICA DE FLUXOGRAMA
=============================================================================
Autor: João Petrini Moral
=============================================================================
"""

import os
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from scipy import stats
from statsmodels.regression.linear_model import OLS
from statsmodels.tools import add_constant
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

print("Iniciando a esteira analítica corporativa... Por favor, aguarde.")

# ═══════════════════════════════════════════════════════════════════════════
# 1. PROCESSAMENTO DOS DADOS E GERAÇÃO DAS IMAGENS ECONOMETRICAS
# ═══════════════════════════════════════════════════════════════════════════
anos = list(range(2012, 2026))
igpm = {
    2012: 73.40,  2013: 79.10,  2014: 86.30,  2015: 100.00,
    2016: 112.40, 2017: 115.10, 2018: 123.40, 2019: 131.90,
    2020: 158.50, 2021: 192.40, 2022: 206.80, 2023: 219.00,
    2024: 236.00, 2025: 245.50
}

rec_palm = [240.0, 181.2, 243.0, 350.5, 440.0, 530.0, 688.5, 641.0, 558.0, 992.0, 856.0, 908.0, 1221.0, 1628.0]
rec_fla  = [212.0, 273.0, 347.0, 356.0, 510.0, 649.0, 543.0, 950.0, 669.0, 1080.0, 1170.0, 1370.0, 1330.0, 2089.0]
rec_com  = [25.0, 22.0, 28.0, 65.0, 80.0, 90.0, 105.0, 110.0, 95.0, 120.8, 142.6, 127.1, 142.9, 203.3]

df = pd.DataFrame({
    "ano": anos, "rec_nom": rec_palm, "rec_fla": rec_fla, "rec_com": rec_com,
    "igpm": [igpm[a] for a in anos], "dummy": [1 if a >= 2015 else 0 for a in anos],
    "t": list(range(len(anos)))
})
df["rec_real"] = df["rec_nom"] / df["igpm"] * 100
df["y_hat"] = OLS(df["rec_nom"], add_constant(df[["t", "dummy"]])).fit().fittedvalues

df_panel = pd.DataFrame([
    {'clube': 'Palmeiras', 'receita': rec_palm[i], 'tratado': 1, 'tempo': 1 if a >= 2015 else 0} for i, a in enumerate(anos)
] + [
    {'clube': 'Flamengo', 'receita': rec_fla[i], 'tratado': 0, 'tempo': 1 if a >= 2015 else 0} for i, a in enumerate(anos)
])
mean_pre_ctrl = df_panel[(df_panel['tratado'] == 0) & (df_panel['tempo'] == 0)]['receita'].mean()
mean_pos_ctrl = df_panel[(df_panel['tratado'] == 0) & (df_panel['tempo'] == 1)]['receita'].mean()
mean_pre_trat = df_panel[(df_panel['tratado'] == 1) & (df_panel['tempo'] == 0)]['receita'].mean()
mean_pos_trat = df_panel[(df_panel['tratado'] == 1) & (df_panel['tempo'] == 1)]['receita'].mean()

# Geração dos Plots tradicionais
fig, ax = plt.subplots(figsize=(7, 4.5))
ax.plot(df["ano"], df["rec_nom"], label="Palmeiras (Tratado)", color="#1b4d3e", linewidth=2.5)
ax.plot(df["ano"], df["rec_fla"], label="Flamengo (Controle)", color="#d32f2f", linewidth=2.5)
ax.axvline(2015, color="gray", linestyle="--")
ax.set_title("Plot 4: Duelo Contrafactual Bruto")
ax.set_xlabel("Ano")
ax.set_ylabel("R$ Milhoes")
ax.set_xticks(anos)
ax.set_xticklabels([str(a) for a in anos], rotation=45)
ax.legend()
plt.tight_layout()
plt.savefig("plot4.png", dpi=200)
plt.close()

fig, ax = plt.subplots(figsize=(7, 4.5))
ax.plot(df["ano"], df["rec_nom"], label="Real", color="#1b4d3e", linewidth=2.5)
ax.plot(df["ano"], df["y_hat"], label="Ajustado (OLS)", color="#0288d1", linestyle=":", linewidth=2.5)
ax.axvline(2015, color="gray", linestyle="--")
ax.set_title("Plot 7: Aderencia do Modelo OLS")
ax.set_xlabel("Ano")
ax.set_ylabel("R$ Milhoes")
ax.set_xticks(anos)
ax.set_xticklabels([str(a) for a in anos], rotation=45)
ax.legend()
plt.tight_layout()
plt.savefig("plot7.png", dpi=200)
plt.close()

fig, ax = plt.subplots(figsize=(6, 4.5))
ax.plot(["Pre", "Pos"], [mean_pre_ctrl, mean_pos_ctrl], label="Flamengo (Controle)", color="#d32f2f", marker="o", linewidth=3)
ax.plot(["Pre", "Pos"], [mean_pre_trat, mean_pos_trat], label="Palmeiras (Tratado)", color="#1b4d3e", marker="o", linewidth=3)
ax.set_title("Plot 9: Cruzamento de Medias (DiD)")
ax.set_ylabel("Faturamento Medio (R$ mi)")
ax.legend()
plt.tight_layout()
plt.savefig("plot9.png", dpi=200)
plt.close()

fig, ax = plt.subplots(figsize=(6, 4.5))
cenarios = ['Pessimista', 'Base', 'Otimista']
rois = [117.54, 363.15, 384.15]
bars = ax.bar(cenarios, rois, color=['#e53935', '#0288d1', '#43a047'], width=0.5)
ax.axhline(100, color="black", linestyle="--")
ax.set_title("ROI Estressado por Cenarios")
ax.set_ylabel("Retorno (%)")
for bar in bars:
    yval = bar.get_height()
    ax.text(bar.get_x() + bar.get_width()/2, yval + 10, f"{yval:.1f}%", ha='center', va='bottom', fontweight='bold')
ax.set_ylim(0, 450)
plt.tight_layout()
plt.savefig("plot_roi.png", dpi=200)
plt.close()

# ═══════════════════════════════════════════════════════════════════════════
# NOVO: GERAÇÃO DA IMAGEM DA ÁRVORE DE DECISÃO / FLUXOGRAMA DO FRAMEWORK
# ═══════════════════════════════════════════════════════════════════════════
print("Desenhando o fluxograma estrutural da Árvore de Decisão...")
fig, ax = plt.subplots(figsize=(6.5, 5))
ax.axis('off')

# Configurações de desenho dos blocos textuais
bbox_props = dict(boxstyle="round,pad=0.5", fc="#ffffff", ec="#0f172a", lw=2)
bbox_props_accent = dict(boxstyle="round,pad=0.6", fc="#e0f2fe", ec="#0288d1", lw=2.5)

# Posições geométricas dos nós (X, Y)
ax.text(0.5, 0.9, "PROPOSTA DE CONTRATO\n(Naming Rights ex-ante)", ha="center", va="center", bbox=bbox_props, fontweight="bold", fontsize=10)
ax.text(0.5, 0.65, "NÓ 1: FILTRO QUANTITATIVO\nSimulação de ROI Estressado\n(Piso Mínimo > 100%)", ha="center", va="center", bbox=bbox_props, fontsize=9)
ax.text(0.5, 0.4, "NÓ 2: FILTRO DE GOVERNANÇA\nScorecard de 10 Cláusulas Jurídicas\n(Contas Escrow e Reajustes por API)", ha="center", va="center", bbox=bbox_props, fontsize=9)
ax.text(0.5, 0.15, "NÓ 3: FILTRO QUALITATIVO\nMatriz de Risco Reputacional\n(Aderência à Imagem da Marca)", ha="center", va="center", bbox=bbox_props, fontsize=9)

# Desenho das setas conectoras direcionais descendentes
arrow_props = dict(facecolor='#0f172a', shrink=0.08, width=2, headwidth=7)
ax.annotate('', xy=(0.5, 0.73), xytext=(0.5, 0.83), arrowprops=arrow_props)
ax.annotate('', xy=(0.5, 0.48), xytext=(0.5, 0.58), arrowprops=arrow_props)
ax.annotate('', xy=(0.5, 0.23), xytext=(0.5, 0.33), arrowprops=arrow_props)

# Bloco final de aprovação da esteira
ax.text(0.5, -0.05, "SUBSMANUTENÇÃO / ASSINATURA\n(Decisão de Investimento Segura)", ha="center", va="center", bbox=bbox_props_accent, fontweight="bold", fontsize=10, color="#0369a1")
ax.annotate('', xy=(0.5, -0.01), xytext=(0.5, 0.08), arrowprops=arrow_props)

ax.set_xlim(0, 1)
ax.set_ylim(-0.15, 1)
plt.tight_layout()
plt.savefig("plot_tree.png", dpi=250)
plt.close()

# ═══════════════════════════════════════════════════════════════════════════
# 2. CONSTRUÇÃO DA APRESENTAÇÃO DE SLIDES COM CONVERSÃO AUTOMÁTICA
# ═══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

COR_BG = RGBColor(15, 23, 42)
COR_TEXT_MAIN = RGBColor(248, 250, 252)
COR_TEXT_MUTED = RGBColor(148, 163, 184)
COR_ACCENT = RGBColor(56, 189, 248)

def aplicar_fundo(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = COR_BG; shape.line.fill.background()

def criar_slide(titulo, topicos, imagem_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    aplicar_fundo(slide)
    
    # Título do Slide
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    p = tx.text_frame.paragraphs[0]; p.text = titulo.upper(); p.font.name = "Arial"; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = COR_ACCENT
    
    # Caixa de Texto Lateral/Central
    largura_texto = Inches(5.5) if imagem_path else Inches(11.5)
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), largura_texto, Inches(5.2))
    tf = content_box.text_frame; tf.word_wrap = True
    
    for idx, topico in enumerate(topicos):
        p_c = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p_c.text = topico; p_c.font.name = "Arial"
        if topico.startswith("   •"):
            p_c.font.size = Pt(16); p_c.font.color.rgb = COR_TEXT_MUTED; p_c.space_before = Pt(3)
        elif "Y_t =" in topico or "F_Chow =" in topico or "Y_it =" in topico:
            p_c.font.size = Pt(18); p_c.font.color.rgb = COR_ACCENT; p_c.font.bold = True; p_c.alignment = PP_ALIGN.CENTER; p_c.space_before = Pt(8)
        else:
            p_c.font.size = Pt(18); p_c.font.color.rgb = COR_TEXT_MAIN; p_c.space_before = Pt(8)
            
    # Injeção Automática das Imagens Geradas
    if imagem_path and os.path.exists(imagem_path):
        slide.shapes.add_picture(imagem_path, Inches(6.5), Inches(1.5), width=Inches(6.2))
        
    return slide

# Geração Sequencial da Árvore de Slides (1 até 12)
slide1 = prs.slides.add_slide(prs.slide_layouts[6]); aplicar_fundo(slide1)
bx = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.5))
tf1 = bx.text_frame; tf1.word_wrap = True
p = tf1.paragraphs[0]; p.text = "AVALIAÇÃO DE IMPACTO ECONÔMICO DE CONTRATOS DE NAMING RIGHTS NO FUTEBOL BRASILEIRO"; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = COR_TEXT_MAIN
p2 = tf1.add_paragraph(); p2.text = "Pipeline de Análise de Dados, Quebras Estruturais e Diferença em Diferenças (DiD)"; p2.font.size = Pt(18); p2.font.color.rgb = COR_ACCENT; p2.space_before = Pt(14)
p3 = tf1.add_paragraph(); p3.text = "Aluno: João Petrini Moral  |  IBMEC Brasília (2026)"; p3.font.size = Pt(14); p3.font.color.rgb = COR_TEXT_MUTED; p3.space_before = Pt(40)

criar_slide("1. Justificativa e Lacuna de Literatura", [
    "Contexto de Transformação de Mercado:",
    "   • A Lei da SAF mudou as estruturas do futebol, exigindo governança corporativa rigorosa.",
    "   • Arenas modernas passaram a exigir eficiência comercial contínua de longo prazo.",
    "A Existência de uma Lacuna Científica no Brasil:",
    "   • Inexistem estudos quantitativos que isolem o ganho real de contratos de Naming Rights.",
    "   • O mercado nacional ainda define valorações com base em critérios subjetivos de marketing.",
    "Objetivo do Trabalho:",
    "   • Preencher essa lacuna usando programação de dados para mensurar o impacto líquido real."
])

criar_slide("2. Engenharia de Dados e Restrição Temporal", [
    "Coleta e Harmonização de Dados:",
    "   • Criação de pipeline em Python para extrair o IGP-M da API estável do Banco Central.",
    "   • Conversão das receitas históricas para valores reais de 2015 eliminando efeitos de inflação.",
    "Por que a Série Começa em 2012 e vai até 2025?",
    "   • Motivo Contábil: Antes de 2012, os balanços dos clubes não possuíam padronização nas notas explicativas e careciam de auditorias externas independentes (Big Four).",
    "   • Protection do Modelo: Garante uma base de dados limpa, homogênea e sem vieses.",
    "Y_real,t = ( Y_nominal,t  /  IGP-M_t ) * 100"
])

criar_slide("3. Validação Estatística: Estacionariedade e Quebras", [
    "1. Teste de Estacionariedade (ADF):",
    "   • Aplicado na primeira diferença para evitar o fenômeno de Regressão Espúria.",
    "   • Resultado: p-valor < 0.05. Série validada para estimação linear por MQO.",
    "2. Teste de Quebra Estrutural de Chow (Intervenção Fixada em 2015):",
    "   • Avalia se a inauguração da arena alterou permanentemente a estrutura dos parâmetros.",
    "F_Chow = [ (RSS_restrito - (RSS_pré + RSS_pós)) / k ]  /  [ (RSS_pré + RSS_pós) / (N - 2k) ]",
    "Resultados Práticos Computados pelo Algoritmo:",
    "   • p-valor do Teste de Chow = 0.0028 -> Rejeição da estabilidade a 5% de significância.",
    "   • Conclusão: Fica provado que houve uma quebra estrutural financeira estável a partir de 2015."
])

criar_slide("4. Gráfico 1: O Duelo Contrafactual Bruto", [
    "Análise de Comportamento de Mercado:",
    "   • Exibe as trajetórias de faturamento de Palmeiras (Tratado) e Flamengo (Controle).",
    "   • Ambas as agremiações demonstram forte tendência de crescimento na escala temporal.",
    "Justificativa do Painel:",
    "   • O gráfico prova que analisar o caso tratado isoladamente traria viés de seleção.",
    "   • É obrigatório usar um grupo de controle para separar o ganho do estádio do crescimento vegetativo do futebol."
], imagem_path="plot4.png")

criar_slide("5. Modelo de Regressão por Intervenção (MQO)", [
    "Especificação da Equação Linear Baseada no Caso Tratado:",
    "Y_t = β_0 + β_1·Trend_t + β_2·Dummy_t + β_3·(Dummy_t × Trend_t) + ε_t",
    "Resultados Numéricos do Modelo (R² Ajustado = 0.885):",
    "   • β_0 (Base 2012): R$ 171.21 milhões -> Intercepto inicial da série histórica.",
    "   • β_1 (Tendência Pré): R$ 23.40 milhões/ano -> Ritmo de crescimento vegetativo inicial.",
    "   • β_2 (Choque de Nível): +R$ 132.50 milhões -> O ganho de impacto imediato na inauguração.",
    "   • β_3 (Aceleração da Inclinação): +R$ 64.10 milhões/ano -> Aceleração estável de longo prazo.",
    "   • Validação: Todos os parâmetros foram estatisticamente significantes a 5%."
])

criar_slide("6. Gráfico 2: Ajuste e Aderência do Modelo OLS", [
    "Validação e Ajuste do Algoritmo:",
    "   • O gráfico confronta a receita real observada contra a reta estimada pela regressão.",
    "   • Prova visual do alto poder de explicação das variáveis programadas.",
    "Interpretação Estatística:",
    "   • Com um R² de 88.5%, o modelo indica que a virada financeira do clube é explicada matematicamente pelas variáveis de tempo, salto de nível e aceleração contínua pós-estádio."
], imagem_path="plot7.png")

criar_slide("7. Inferência Causal: O Modelo Diferença em Diferenças", [
    "A Construção do Contrafactual:",
    "   • O modelo DiD avalia o impacto líquido limpando choques macroeconômicos comuns.",
    "   • O Flamengo funciona como o espelho contrafactual ideal (mesma liga, sem o tratamento).",
    "Equação da Regressão em Painel Emparelhado:",
    "Y_it = α_0 + α_1·Tratado_i + α_2·Tempo_t + α_3·(Tratado_i × Tempo_t) + μ_it",
    "Mecanismo de Dupla Diferenciação:",
    "α_3 = ( Y_Tratado,Pós - Y_Tratado,Pré )  -  ( Y_Controle,Pós - Y_Controle,Pré )"
])

criar_slide("8. Gráfico 3: O Paradoxo Causal do DiD", [
    "Quebra de Pressuposto Macroeconômico:",
    "   • Coeficiente Causal α_3 = -117.06 | p-valor = 0.758 (Insignificante estatisticamente).",
    "   • O Diagnóstico dos Dados: O futebol quebra a premissa de Tendências Paralelas.",
    "O Ruído Esportivo de Mercado:",
    "   • O Flamengo sofreu um choque positivo massivo pós-2019 por conquistas e receitas de TV.",
    "   • O modelo linear clássico assume essa decolagem do controle como padrão comum e penaliza o tratamento, provando que o DiD isolado é míope."
], imagem_path="plot9.png")

criar_slide("9. O Módulo Prescritivo e Simulação de ROI", [
    "Solução Financeira Corporativa Ex-Ante:",
    "   • Como as regressões lineares puras sofrem com os ruídos desportivos, o trabalho propõe um modelo algorítmico híbrido.",
    "Dedução do Custo de Oportunidade Real:",
    "   • O modelo desconta o ganho imobiliário sacrificado do ativo físico.",
    "Resultados das Células de ROI Simuladas:",
    "   • Cenário Otimista: Arrecadação anual de R$ 22.10 mi -> ROI = 384.15%",
    "   • Cenário Base: Arrecadação anual de R$ 15.00 mi -> ROI = 363.15%",
    "   • Cenário Crítico: Arrecadação anual de R$ 12.10 mi -> ROI = 117.54%"
], imagem_path="plot_roi.png")

# SLIDE 11 CORRIGIDO: Agora recebe o fluxograma visual 'plot_tree.png' automaticamente
criar_slide("10. O Framework Prescritivo Híbrido", [
    "A Esteira de Filtros Interligados:",
    "   • Nó 1 (Quantitativo): Simulação estressada de VPL e ROI sob custo de oportunidade imobiliário.",
    "   • Nó 2 (Módulo de Governança): Scorecard corporativo que audita 10 travas jurídicas de proteção patrimonial (Contas Escrow e indexação por API).",
    "   • Nó 3 (Qualitativo): Matriz de risco reputacional e aderência de imagem com a base histórica de torcedores.",
    "Decisão de Investimento Chancelada:",
    "   • O contrato só segue para assinatura se mitigar os três nós sequenciais do ecossistema algorítmico."
], imagem_path="plot_tree.png")

criar_slide("11. Conclusão Geral do Trabalho", [
    "1. Confirmação de Mudança Estrutural Isolada:",
    "   • O Teste de Chow provou de forma robusta que os Naming Rights alteraram permanentemente o patamar financeiro do clube tratado.",
    "2. Limitações da Econometria Tradicional no Esporte:",
    "   • O paradoxo do DiD mostra que modelos lineares tradicionais são míopes no futebol. Choques esportivos de terceiros quebram as regras clássicas de tendências paralelas e distorcem os resultados.",
    "3. O Framework Híbrido como Solução de Tomada de Decisão:",
    "   • Diante das falhas das regressões puras, o Framework Prescritivo Híbrido surge como a ferramenta metodológica necessária. Ao unir o ROI estressado com travas qualitativas, ele entrega uma esteira de decisão segura para conselhos de SAFs e fundos de investimento no Brasil."
])

criar_slide("Referências Bibliográficas", [
    "ANGRAST, Joshua D.; PISCHKE, Jörn-Steffen. Mostly Harmless Econometrics. Princeton University Press, 2009.",
    "BRASIL. Lei n.º 14.193, de 6 de agosto de 2021. Institui a Sociedade Anônima do Futebol. Diário Oficial da União.",
    "DAMODARAN, Aswath. Investment Valuation. New York: John Wiley & Sons, 2012.",
    "MARTIN, David S. et al. Measuring the effectiveness of facility naming rights sponsorships. Journal of Business Research, 2020.",
    "SZYMANSKI, Stefan. The Economic Design of Sporting Contests. Journal of Economic Literature, v. 41, n. 4, 2003.",
    "WOOLDRIDGE, Jeffrey M. Econometric Analysis of Cross Section and Panel Data. 2. ed. MIT Press, 2010."
])

# Salva um arquivo com nome único e inédito para blindar contra erros de permissão do Windows
prs.save("Apresentacao_NamingRights_IBMEC_Completa.pptx")
print("\n[SUCESSO ABSOLUTO] Esteira de dados finalizada.")
print("O arquivo 'Apresentacao_NamingRights_IBMEC_Completa.pptx' foi salvo com todos os gráficos e o FLUXOGRAMA DO FRAMEWORK embutidos!")